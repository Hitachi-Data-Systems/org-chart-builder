from collections import OrderedDict
import math
import datetime
import pprint
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import re
from shape import RectangleBuilder, ColorPicker

__author__ = 'David Oreper'


class DrawChartSlide:
    MAX_WIDTH_INCHES = Inches(10)
    MAX_HEIGHT_INCHES = Inches(5.63)
    PAGE_BUFFER = Inches(.2)
    RIGHT_EDGE = MAX_WIDTH_INCHES - PAGE_BUFFER
    FOOTER_TOP = MAX_HEIGHT_INCHES - Inches(.45)
    FOOTER_LEFT = Inches(.7)


    def __init__(self, aPresentation, slideTitle, slideLayout, teamModelText=None):
        """

        :type aPresentation: pptx.api.Presentation
        :type slideTitle: str
        :type slideLayout: pptx.parts.slidelayout.SlideLayout
        """
        self.presentation = aPresentation
        self.slideTitle = slideTitle
        self.slideLayout = slideLayout
        self.colorPicker = ColorPicker()
        self.groupList = []
        self.groupLeft = Inches(.2)
        self.groupTop = GroupShapeDimensions.TOP
        self.teamModelText = teamModelText

    def addGroup(self, title, groupMembers):
        peopleGroup = self._getPeopleGroup(title)
        peopleGroup.setGroupColor(self.colorPicker.getBackgroundColor())
        peopleGroup.setMemberColor(self.colorPicker.getForegroundColor())
        self.colorPicker.nextColor()

        peopleGroup.addMembers(groupMembers)

        self.groupLeft = peopleGroup.getRightEdge() + GroupShapeDimensions.BUFFER_WIDTH
        self.groupList.append(peopleGroup)

    def _getPeopleGroup(self, title):
        return PeopleGroup(title, self.groupLeft, GroupShapeDimensions.TOP, GroupShapeDimensions.HEIGHT)

    def _adjustGroupWidth(self, reductionRatio):
        print("Reducing width by {}% for: {}".format(100 - int(reductionRatio * 100), self.slideTitle))
        rightEdge = Inches(.2)
        for aGroup in self.groupList:
            aGroup.adjustWidth(reductionRatio)
            aGroup.setLeft(rightEdge)
            rightEdge = aGroup.getRightEdge() + GroupShapeDimensions.BUFFER_WIDTH * reductionRatio

    def _center(self):
        rightEdge = self.groupList[-1].getRightEdge()
        leftAdjust = (DrawChartSlide.RIGHT_EDGE - rightEdge) / 2
        for aGroup in self.groupList:
            aGroup.setLeft(aGroup.getLeft() + leftAdjust)

    def addTitle(self, aSlide):

        aSlide.shapes.title.text = self.slideTitle

        modelText = ""
        if self.teamModelText:
            modelText = "\tModel: {}".format(self.teamModelText)

        effectiveDateRun = aSlide.shapes.title.textframe.paragraphs[0].add_run()

        # A little ugly but does the trick - add suffix to the number: 2 -> 2nd
        # http://stackoverflow.com/a/20007730
        suffixFormatter = lambda n: "%d%s" % (n, "tsnrhtdd"[(n / 10 % 10 != 1) * (n % 10 < 4) * n % 10::4])
        ordinalDay = suffixFormatter(int(datetime.datetime.now().strftime("%d").lstrip("0")))

        month = datetime.datetime.now().strftime("%B")
        effectiveDateRun.text = "\rEffective {} {} {}".format(month, ordinalDay, modelText)
        effectiveDateRun.font.size = Pt(11)
        effectiveDateRun.font.italic = True
        effectiveDateRun.font.bold = False

    def addFooter(self, aSlide, footerString, leftEdge=0):
        footerTextBox = aSlide.shapes.add_textbox(leftEdge, DrawChartSlide.FOOTER_TOP, Inches(0), Inches(0))
        footerTextFrame = footerTextBox.textframe
        p = footerTextFrame.add_paragraph()
        p.text = footerString
        p.font.size = Pt(7)


    def drawSlide(self):
        if not self.groupList:
            print "WARNING: NO Groups added for product: {}".format(self.slideTitle)
            return

        slide = self.presentation.slides.add_slide(self.slideLayout)
        self.addTitle(slide)

        rightEdge = self.groupList[-1].getRightEdge()
        if rightEdge > DrawChartSlide.RIGHT_EDGE:
            reductionRatio = 1 / (float(rightEdge) / DrawChartSlide.RIGHT_EDGE)
            self._adjustGroupWidth(reductionRatio)

        self._center()

        totalMembers = 0
        totalTBH = 0
        totalExpat = 0
        for aGroup in self.groupList:
            aGroup.build(slide)
            totalMembers += len(aGroup.memberShapeList)- aGroup.totalExpat - aGroup.totalTBH
            totalTBH += aGroup.totalTBH
            totalExpat += aGroup.totalExpat

        footerString = ""
        if totalMembers:
            footerString = "|| HC:{}  ".format(totalMembers)
        if totalTBH:
            footerString += "||  TBH:{}  ".format(totalTBH)
        if totalExpat:
            footerString += "||  Expat:{}  ".format(totalExpat)

        if footerString:
            footerString += "||"
        self.addFooter(slide, footerString, self.groupList[0].groupLeft)


class DrawChartSlideAdmin(DrawChartSlide):
    def _getPeopleGroup(self, title):
        return PeopleGroupAdmin(title, self.groupLeft, GroupShapeDimensions.TOP, GroupShapeDimensions.HEIGHT)


class DrawChartSlideTBH(DrawChartSlide):
    def _getPeopleGroup(self, title):
        return PeopleGroupTBH(title, self.groupLeft, GroupShapeDimensions.TOP, GroupShapeDimensions.HEIGHT)

class DrawChartSlideExpatIntern(DrawChartSlide):
    def _getPeopleGroup(self, title):
        return PeopleGroupExpatIntern(title, self.groupLeft, GroupShapeDimensions.TOP, GroupShapeDimensions.HEIGHT)

class GroupShapeDimensions:
    def __init__(self):
        pass

    TOP = Inches(1.1)
    HEIGHT = Inches(4.3)
    WIDTH = Inches(1.17)
    BUFFER_WIDTH = Inches(.05)
    BUFFER_HEIGHT = Inches(.4)


class MemberShapeDimensions:
    def __init__(self):
        pass

    HEIGHT = Inches(.45)
    BUFFER_WIDTH = Inches(.03)
    BUFFER_HEIGHT = Inches(.03)
    WIDTH = GroupShapeDimensions.WIDTH - (BUFFER_WIDTH * 2)
    HARD_WRAP_NUM = (GroupShapeDimensions.HEIGHT - GroupShapeDimensions.BUFFER_HEIGHT) / (HEIGHT + BUFFER_HEIGHT)

class PeopleGroup(object):
    def __init__(self, title, left, top, height):
        """

        :type title: str
        :type left: Inches or float
        :type top: pptx.util.Inches
        :type height: pptx.util.Inches
        """
        self.name = title
        self.title = title
        self.groupWidthUnits = 1
        self.groupUnitWidth = GroupShapeDimensions.WIDTH
        self.groupTop = top
        self.groupLeft = left
        self.groupHeight = height
        self.totalTBH = 0
        self.totalExpat = 0

        self.memberLeft = self.groupLeft + MemberShapeDimensions.BUFFER_WIDTH
        self.memberTop = self.groupTop + GroupShapeDimensions.BUFFER_HEIGHT
        self.memberShapeList = []
        self.fontReduce = 1

    def _nextColumn(self):
        self.memberTop = self.groupTop + GroupShapeDimensions.BUFFER_HEIGHT
        self.memberLeft = self.memberLeft + MemberShapeDimensions.WIDTH + MemberShapeDimensions.BUFFER_WIDTH * 2

    def setGroupColor(self, groupColorRGB):
        """

        :type groupColorRGB: pptx.dml.color.RGBColor
        """
        self.groupColorRGB = groupColorRGB

    def setMemberColor(self, memberColorRGB):
        """

        :type memberColorRGB: pptx.dml.color.RGBColor
        """
        self.memberColor = memberColorRGB

    def setHeight(self, heightInches):
        """

        :param heightInches:
        """
        self.groupHeight = heightInches

    def getLeft(self):
        return self.groupLeft

    def setLeft(self, leftCoord):
        """

        :type leftCoord: float or Inches
        """
        self.groupLeft = leftCoord

        leftStart = self.memberShapeList[0].getLeft()
        memberLeft = self.groupLeft + MemberShapeDimensions.BUFFER_WIDTH
        for aMember in self.memberShapeList:
            if (aMember.getLeft() - leftStart) > 0:
                leftStart = aMember.getLeft()
                memberLeft += aMember.getWidth() + MemberShapeDimensions.BUFFER_WIDTH

            aMember.setLeft(memberLeft)

    def adjustWidth(self, reductionRatio):
        """

        :type reductionRatio: float
        """
        self.groupUnitWidth = self.groupUnitWidth * reductionRatio
        self.fontReduce = reductionRatio + ((1 - reductionRatio) / 2)
        for aMember in self.memberShapeList:
            aMember.adjustWidth(reductionRatio)
            aMember.adjustFontSizes(self.fontReduce)

    def getWidth(self):
        return self.groupWidthUnits * self.groupUnitWidth

    def getRightEdge(self):
        return self.groupLeft + self.getWidth()

    def addMembers(self, peopleList):
        count = 1

        # Calculate the total width of the group so that we can distribute columns evenly
        self.groupWidthUnits = max(math.ceil(len(peopleList) / float(MemberShapeDimensions.HARD_WRAP_NUM)), 1)

        # Calculate how many people will be in each column
        wrapCount = math.ceil(len(peopleList) / float(self.groupWidthUnits))
        for aPerson in peopleList:

            aPersonRect = self._getPersonRect(aPerson)
            aPersonRect = self.addRectFormatting(aPerson, aPersonRect)
            self.addMemberRect(aPersonRect)

            count += 1
            if count > wrapCount:
                self._nextColumn()
                count = 1

    def _getPersonRect(self, aPerson):
        """

        :param aPerson:
        :return: Return a personRect with FirstName,LastName,Title pre-populated
        """
        aPersonRect = RectangleBuilder(self.memberLeft, self.memberTop, MemberShapeDimensions.WIDTH,
                                       MemberShapeDimensions.HEIGHT)

        firstName = aPerson.getFirstName()
        lastName = aPerson.getLastName()
        if aPerson.isTBH() and ("(" in firstName):
            lastName = "(" + "(".join(firstName.split("(")[1:])
            firstName = firstName.split("(")[0]

        title = self._getTitle(aPerson)

        aPersonRect.setFirstName(firstName)
        aPersonRect.setLastName(lastName)
        aPersonRect.setTitle(title)
        return aPersonRect

    def _getTitle(self, aPerson):

        if aPerson.isTBH():
            if re.search('\d\d\d\d\d', aPerson.getRawName()):
                return aPerson.getRawName()
            return "{}\n{}".format(aPerson.getTitle(), aPerson.getReqNumber())
        else:
            if aPerson.isConsultant():
                return aPerson.getTitle() + " (c)"
            if aPerson.isVendor():
                return aPerson.getTitle() + " (v)"
            if aPerson.isExpat():
                return aPerson.getTitle() + " (e)"
            if aPerson.isIntern():
                return aPerson.getTitle() + " (i)"
            else:
                return aPerson.getTitle()

    def addRectFormatting(self, aPerson, aPersonRect):
        aPersonRect.setBrightness(0)
        aPersonRect.setRGBFillColor(self.memberColor)
        aPersonRect.setRGBTextColor(RGBColor(255, 255, 255))
        aPersonRect.setRGBFirstNameColor(RGBColor(255, 255, 255))

        if aPerson.isLead():
            aPersonRect.setRGBFirstNameColor(RGBColor(127, 127, 127))

        if aPerson.isManager():
            aPersonRect.setRGBFirstNameColor(RGBColor(255, 238, 0))

        if aPerson.isTBH():
            self.totalTBH += 1
            if self.isFutureTBH(aPerson):
                aPersonRect.setBrightness(.4)

        if aPerson.isExpat():
            self.totalExpat += 1

        return aPersonRect

    def addMemberRect(self, aPersonRect):
        self.memberShapeList.append(aPersonRect)
        self.memberTop += MemberShapeDimensions.HEIGHT + MemberShapeDimensions.BUFFER_HEIGHT

    def build(self, aSlide):
        groupRect = RectangleBuilder(self.groupLeft, self.groupTop, self.getWidth(), self.groupHeight)
        groupRect.setRGBFillColor(self.groupColorRGB)
        groupRect.setRGBTextColor(RGBColor(0, 0, 0))
        groupRect.setHeading(self.title)
        groupRect.setBrightness(.2)
        groupRect.adjustFontSizes(self.fontReduce)
        groupRect.build(aSlide)
        for aMember in self.memberShapeList:
            aMember.build(aSlide)

        self.addGroupFooter(aSlide)


    def addGroupFooter(self, aSlide):
        footerTextBox = aSlide.shapes.add_textbox(self.groupLeft - GroupShapeDimensions.BUFFER_WIDTH,
                                                  GroupShapeDimensions.HEIGHT + GroupShapeDimensions.TOP
                                                  - GroupShapeDimensions.BUFFER_HEIGHT
                                                  + MemberShapeDimensions.BUFFER_HEIGHT ,
                                                  Inches(.5), Inches(0))
        footerTextFrame = footerTextBox.textframe
        p = footerTextFrame.add_paragraph()
        p.text = pprint.pformat(len(self.memberShapeList))
        p.font.size = Pt(5)
        p.font.italic = True

    def isFutureTBH(self, aPerson):
        return datetime.datetime.now() < aPerson.getStartDate()

class PeopleGroupAdmin(PeopleGroup):
    def _getInternExpatTitle(self, aPerson):
        return aPerson.getTitle()

class PeopleGroupTBH(PeopleGroup):
    def _getPersonRect(self, aPerson):
        aPersonRect = RectangleBuilder(self.memberLeft, self.memberTop, MemberShapeDimensions.WIDTH,
                                       MemberShapeDimensions.HEIGHT)

        firstName = aPerson.getFunction()
        title = "{}".format(aPerson.getTitle())
        if aPerson.isTBH() and ("(" in aPerson.getFirstName()):
            title = "{} ({}".format(title, "(".join(aPerson.getFirstName().split("(")[1:]))

        if self.isFutureTBH(aPerson):
            # Use the person's name instead of the req number if this is a future hire
            reqNumber = aPerson.getFullName()
        else:
            reqNumber = aPerson.getReqNumber()
        
        aPersonRect.setFirstName(firstName)
        aPersonRect.firstNameSize = 7

        # If title is too long, only show the req number
        if len(aPerson.getTitle()) + len(aPerson.getFunction()) > 37:
            aPersonRect.setLastName(reqNumber)
        else:
            aPersonRect.setLastName(title)
            aPersonRect.setTitle(reqNumber)
        return aPersonRect

class PeopleGroupExpatIntern(PeopleGroup):
    def _getPersonRect(self, aPerson):
        aPersonRect = RectangleBuilder(self.memberLeft, self.memberTop, MemberShapeDimensions.WIDTH,
                                       MemberShapeDimensions.HEIGHT)

        aPersonRect.setFirstName(aPerson.getFirstName())
        aPersonRect.setLastName(aPerson.getLastName())
        aPersonRect.setTitle(aPerson.getProduct())
        return aPersonRect
