from pptx import Presentation
from pptx.dml.color import RGBColor

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

__author__ = 'David Oreper'

class ColorPicker:
    def __init__(self):
        self.index = 0
        self.colors = [(RGBColor(245, 227, 226), RGBColor(166, 75, 71)),
                       (RGBColor(209, 225, 243), RGBColor(42, 94, 142)),
                       (RGBColor(239, 243, 229), RGBColor(138, 160, 81)),
                       (RGBColor(227, 235, 244), RGBColor(72, 117, 161)),
                       (RGBColor(235, 231, 240), RGBColor(116, 96, 140)),
                       (RGBColor(226, 241, 246), RGBColor(65, 151, 171)),
                       (RGBColor(253, 238, 226), RGBColor(235, 128, 35))]

    def getBackgroundColor(self):
        colorIndex = self.index % len(self.colors)
        return self.colors[colorIndex][0]

    def getForegroundColor(self):
        colorIndex = self.index % len(self.colors)
        return self.colors[colorIndex][1]

    def nextColor(self):
        self.index += 1


class RectangleBuilder(object):
    """
    Build rectangles
    """

    def __init__(self, left, top, width, height):
        """
        :param left:
        :param top:
        :param width:
        :param height:
        """
        self.width = width
        self.height = height
        self.top = top
        self.left = left
        self.rgbTextColor = RGBColor(0, 0, 0)
        self.rgbFillColor = RGBColor(255, 255, 255)
        self.rgbFirstNameColor = self.rgbTextColor
        self.brightness = 0
        self.firstName = None
        self.nickName = None
        self.lastName = None
        self.title = None
        self.heading = None
        self.firstNameSize = 8
        self.headingSize = 9
        self.lastNameSize = 6
        self.titleSize = 5
        self.minFontSize = 5

    def adjustFont(self, currentSize, scaleRatio):
        return max(self.minFontSize, float(currentSize) * scaleRatio)

    def adjustFontSizes(self, scaleRatio):
        self.headingSize = self.adjustFont(self.headingSize, scaleRatio)
        self.firstNameSize = self.adjustFont(self.firstNameSize, scaleRatio)
        self.lastNameSize = self.adjustFont(self.lastNameSize, scaleRatio)
        self.titleSize = self.adjustFont(self.titleSize, scaleRatio)

    def setLeft(self, leftCoord):
        self.left = leftCoord

    def setTop(self, topCoord):
        self.top = topCoord

    def getTop(self):
        return self.top

    def getWidth(self):
        return self.width

    def getLeft(self):
        return self.left

    def adjustWidth(self, widthAdjustmentRatio):
        self.width = self.width * widthAdjustmentRatio

    def setRGBTextColor(self, rgbColor):
        """

        :type rgbColor: pptx.dml.color.RGBColor
        """
        self.rgbTextColor = rgbColor

    def setRGBFirstNameColor(self, rgbColor):
        self.rgbFirstNameColor = rgbColor

    def setRGBFillColor(self, rgbColor):
        self.rgbFillColor = rgbColor

    def setFirstName(self, firstName):
        """

        :type firstName: str
        """
        self.firstName = firstName

    def setNickName(self, nickName):
        self.nickName = nickName

    def setLastName(self, lastName):
        """

        :type lastName: list or str
        """
        self.lastName = lastName

    def setTitle(self, title):
        """

        :type title: str
        """
        self.title = title

    def setHeading(self, heading):
        """

        :type heading: str
        """
        self.heading = heading

    def setBrightness(self, brightness):
        """

        :type brightness: float or int
        """
        self.brightness = brightness

    def _buildShape(self, aSlide):
        shape = aSlide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.left, self.top, self.width, self.height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.rgbFillColor
        shape.fill.fore_color.brightness = self.brightness
        shape.line.color.rgb = self.rgbFillColor
        return shape

    def addRun(self, textFrame, text, size, isBold, isItalic, rgbTextColor=None):
        """

        :param textFrame:
        :type text: str
        :type size: int
        :type isBold: bool
        :type isItalic: bool
        :type rgbTextColor: __builtin__.NoneType
        """
        paragraph = textFrame.paragraphs[0]
        aRun = paragraph.add_run()

        if not rgbTextColor:
            rgbTextColor = self.rgbTextColor

        text = "{}\r".format(text.strip())
        aRun.text = text

        aRun.font.size = Pt(size)
        aRun.font.bold = isBold
        aRun.font.italic = isItalic
        aRun.font.color.rgb = rgbTextColor
        aRun.alignment = PP_ALIGN.CENTER
        aRun.font.name = "Arial (Body)"

    def _buildFirstName(self, textFrame):
        self.addRun(textFrame, self.firstName, self.firstNameSize, True, False, self.rgbFirstNameColor)

    def _buildLastName(self, textFrame):
        name = self.lastName
        if self.nickName:
            name = "({}){}".format(self.nickName, name)

        self.addRun(textFrame, name, self.lastNameSize, False, False)

    def _buildTitle(self, textFrame):
        title = self.title
        self.addRun(textFrame, title, self.titleSize, False, True)

    def _buildHeading(self, textFrame):
        textFrame.margin_left = 0
        textFrame.margin_right = 0
        self.addRun(textFrame, self.heading, self.headingSize, True, False)

    def build(self, aSlide):
        shape = self._buildShape(aSlide)
        shape.textframe.vertical_anchor = MSO_ANCHOR.TOP

        if self.heading:
            self._buildHeading(shape.textframe)

        if self.firstName:
            self._buildFirstName(shape.textframe)

        if self.lastName:
            self._buildLastName(shape.textframe)

        if self.title:
            self._buildTitle(shape.textframe)


class SlideTitleShape:
    def __init__(self):
        self.top = 0
        self.left = 0
        self.titleSize = 25
        self.width = Inches(10)
        self.height = Inches(1)

    def setTop(self, top):
        self.top = top

    def setLeft(self, left):
        self.left = left

    def setWidth(self, width):
        self.width = width

    def setHeight(self, height):
        self.height = height

    def setTitleSize(self, size):
        self.titleSize = size

    def drawTitle(self, slideTitle, aSlide):
        textBox = aSlide.shapes.add_textbox(self.left, self.top, self.width, self.height)
        paragraph = textBox.textframe.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        aRun = paragraph.add_run()
        text = "{}\r".format(slideTitle.strip())
        aRun.text = text
        aRun.font.size = Pt(self.titleSize)
        aRun.font.bold = True
        aRun.alignment = PP_ALIGN.CENTER
        aRun.font.name = "Arial (Body)"

#
# if __name__ == "__main__":
#     presentation = Presentation()
#     presentation.slide_height = Inches(10)
#     presentation.slide_width = Inches(10)
#
#     slide = presentation.slides.add_slide(presentation.slide_layouts[6])
#     SlideTitleShape("TEST", slide)
#
#     presentation.save("test.pptx")

